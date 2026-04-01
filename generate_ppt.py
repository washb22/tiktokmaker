import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

pptx_path = r"C:\Users\ghtes\Desktop\tiktok-generator\영상분석_보고서.pptx"

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

def add_title_slide(title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.33), Inches(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(20)
        p2.font.color.rgb = RGBColor(0xAA, 0xAA, 0xCC)
        p2.alignment = PP_ALIGN.CENTER

def add_content_slide(title, content_lines, hc=RGBColor(0x00, 0x96, 0xD6)):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(1.1))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = hc
    title_shape.line.fill.background()
    tf = title_shape.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.5)
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12.33), Inches(5.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(content_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        is_header = line.startswith("##")
        is_bullet = line.startswith("-") or line.startswith("*")
        clean = line.lstrip("#-* ")
        p.text = clean
        if is_header:
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
            p.space_before = Pt(12)
        elif is_bullet:
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
            p.level = 1
        else:
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

# 1. Title
add_title_slide(
    "틱톡 뷰티 숏폼 25개 영상 정밀 분석",
    "2026.04.01 | 분석 영상: 15개 신규 + 10개 기존 | 재원 (데이터 분석)"
)

# 2. Overview
add_content_slide("분석 개요", [
    "## 분석 대상",
    "- 전체 25개 틱톡 뷰티/자기관리 숏폼 영상",
    "- 신규 15개 (snaptik 다운로드) + 기존 10개 (한글 파일명)",
    "- 해상도: 720x1280 (15개) / 1080x1920 (10개)",
    "",
    "## 영상 길이 통계",
    "- 전체 25개 평균: 43.3초",
    "- 신규 15개 평균: 53.4초 (최소 30초 ~ 최대 82.6초)",
    "- 기존 10개 평균: 33.1초 (최소 19.2초 ~ 최대 57.5초)",
    "- 최적 범위: 스토리형 50~70초 / 리스트형 30~40초",
    "",
    "## 주제 분포",
    "- 가슴 크기 관련: 32% (8개) - 석류정 제품",
    "- 피부 미백/착색: 28% (7개) - 워시비 착색크림",
    "- 뷰티 루틴: 20% (5개)",
    "- 체형 관리: 12% (3개)",
    "- 기타(키 성장): 4% (1개)"
])

# 3. Hooking
add_content_slide("훅킹 패턴 분석 (핵심 발견)", [
    "## 커뮤니티 게시글 위장 훅킹 - 68% (17/25개)",
    "- 네이트판 '톡커들의선택' 뱃지 + 조회수/댓글수 표시",
    "- 패완표 '오늘의판' 뱃지도 사용",
    "- 평균 표시 조회수: 148,500 / 평균 표시 댓글수: 255개",
    "- '이미 검증된 정보'라는 사회적 증거가 스크롤 멈추게 함",
    "",
    "## 이미지+텍스트 훅킹 - 16% (4개)",
    "- 민감한 주제: 손그림/일러스트로 대체",
    "- 밈 캐릭터(인어공주) 또는 충격 비주얼(팩 바른 얼굴)",
    "",
    "## 텍스트 온리 훅킹 - 12% (3개)",
    "- 큰 볼드체 한마디: '이거봐'",
    "- 스토리 시작형: '내가 몇개월 전만 하더라도'",
    "",
    "## 긴급감 훅킹 - 4% (1개)",
    "- '+1시간뒤 삭제' 같은 삭제 예고"
], RGBColor(0xE0, 0x4D, 0x4D))

# 4. Structure
add_content_slide("장면 구조 유형 4가지", [
    "## A. 스토리 드리븐 (56%, 14개) - 평균 58초, 8장면",
    "- 훅킹 -> 공감스토리 -> 시도/실패 -> 원인설명 -> 솔루션 -> 증거 -> 팁 -> 마무리",
    "- 감정 낙차가 핵심. 좌절 구간이 길수록 솔루션 가치 상승",
    "",
    "## B. 리스트 정보형 (24%, 6개) - 평균 37초, 6장면",
    "- 훅킹 -> 번호리스트(1.2.3.) -> 제품삽입 -> 마무리",
    "- 짧고 빠른 정보 전달. 댓글 참여도 가장 높음 (평균 392)",
    "",
    "## C. 3자 시점형 (12%, 3개) - 평균 37초, 6장면",
    "- 훅킹 -> 친구/주변인 이야기 -> 솔루션 -> 증거 -> 반전 -> 마무리",
    "- '친구가~', '주변에서~' 3인칭으로 광고감 최소화",
    "",
    "## D. 결과 쇼케이스형 (8%, 2개) - 평균 53초, 7장면",
    "- 훅킹 -> 변화리스트 -> 솔루션 -> 증거사진 -> 공유유도 -> 마무리",
    "- 결과부터 보여주고 방법 설명"
], RGBColor(0x00, 0x7A, 0x33))

# 5. Emotion
add_content_slide("감정 곡선 패턴", [
    "## 스토리형 감정 곡선 (56%)",
    "- 호기심(100%) -> 공감(80%) -> 불안(60%) -> 좌절(20%) -> 납득(50%) -> 기대(75%) -> 확신(95%)",
    "",
    "## 핵심 포인트",
    "- 최저점: 시도/실패 후반 (영상 40~50% 지점)",
    "- 전환점: 원인설명 시작 (영상 50~60% 지점)",
    "- 최고점: 증거/후기 (영상 75~85% 지점)",
    "- 40%의 영상에서 마지막에 긴급감 추가",
    "",
    "## 리스트형 감정 곡선 (24%)",
    "- 호기심 -> 납득 -> 기대 -> 확신 (하락 없이 상승형)",
    "",
    "## 3자 시점형 감정 곡선 (12%)",
    "- 호기심 -> 놀라움 -> 기대 -> 확신 -> 공유욕구"
], RGBColor(0x9C, 0x27, 0xB0))

# 6. Image Strategy
add_content_slide("이미지 전략 분석", [
    "## 영상당 평균 이미지 3.4장 / 텍스트 온리 3.8장 (비율 47:53)",
    "",
    "## 이미지 유형별 사용률",
    "- 제품 실물 사진: 92% (솔루션 장면)",
    "- 커뮤니티 스크린샷: 68% (첫 장면 훅킹)",
    "- 신체 사진: 48% (증거/후기 장면)",
    "- 해부도/성분표: 20% (원인설명 장면)",
    "- 밈/캐릭터: 16% (훅킹/리듬감)",
    "- 전후 비교: 12% (증거 장면)",
    "- X표시 실패사진: 8% (시도/실패 장면) - 새 발견!",
    "- 리뷰 캡쳐: 8% (증거 장면)",
    "",
    "## 핵심 인사이트",
    "- 제품 사진은 거의 100% 필수. 자연스러운 실물 사진",
    "- 해부도/성분표가 정보 장면에서 신뢰도 상승",
    "- X표시 실패사진: 기존 제품에 빨간 X 그어 실패 시각화"
], RGBColor(0xFF, 0x98, 0x00))

# 7. Empathy
add_content_slide("공감 유도 방식 Top 5", [
    "## 1. 본인 경험담 고백체 (56%)",
    "- '나 어릴때부터 까맸음', '은근 계속 신경쓰이더라구요'",
    "",
    "## 2. 구체적 실패 나열 (56%)",
    "- '미백크림도 써보고, 병원도 알아봤는데, 수백만원 써도~'",
    "- 최소 3가지 시도+실패 필수. 기간/금액 구체적으로",
    "",
    "## 3. 사회적 압박 언급 (48%)",
    "- '남자들이 ~하면 ~로 본다는 얘기를 듣고'",
    "",
    "## 4. 감정의 행동화 (72%)",
    "- '신경쓰이다'(X) -> '거울 볼때마다 한숨'(O)",
    "",
    "## 5. 3자 시점 이야기 (12%)",
    "- '친구가 방학동안 가슴이 커졌는데 대체 어떻게 한거냐고'"
], RGBColor(0x00, 0x96, 0xD6))

# 8. New Patterns
add_content_slide("새로 발견한 6가지 패턴", [
    "## 1. 커뮤니티 게시글 위장 훅킹 (68%)",
    "- 네이트판/패완표 스크린샷 + 조회수+댓글수 = 사회적 증거",
    "",
    "## 2. 석류정 제품군 집중 포맷 (32%)",
    "- 같은 제품이라도 주제/각도를 바꿔 다양한 영상 생산",
    "",
    "## 3. 짧은 리스트형 30~40초 포맷 (24%)",
    "- 기존 style_profile에 없던 새 구조. 완시청율+댓글 높음",
    "",
    "## 4. 3자 시점 스토리 (12%)",
    "- 친구/주변인 이야기로 광고감 최소화",
    "",
    "## 5. 가짜후기 반박 신뢰 전략",
    "- '100% 좋은 후기만 있는 제품은 없다??' 먼저 의심 던지기",
    "",
    "## 6. 긴급감/희소성 CTA (40%)",
    "- '재고 얼마 없다고', '+1시간뒤 삭제', '지금 시기 놓치면'"
], RGBColor(0xE0, 0x4D, 0x4D))

# 9. Engagement
add_content_slide("참여도 높은 영상 특징", [
    "## 댓글수 TOP 4",
    "- 1위: 다리 이뻐지는법 (댓글 428, 조회 258,927) - 리스트형, 34초",
    "- 2위: 이뻐지는법 (댓글 384, 조회 204,072) - 리스트형, 35초",
    "- 3위: 가슴커지는법 (댓글 384, 조회 147,938) - 스토리형, 69초",
    "- 4위: 가슴커진애 (댓글 364, 조회 93,466) - 3자시점, 39초",
    "",
    "## 공통점",
    "- 리스트형(30~40초)이 댓글 유도에 가장 효과적",
    "- 실용적 정보 + 짧은 길이 = 높은 완시청율 + 댓글 참여",
    "",
    "## 조회수 TOP 3",
    "- 1위: 다리 이뻐지는법 (조회 258,927)",
    "- 2위: 이뻐지는법 (조회 204,072)",
    "- 3위: 가슴 작다면 (조회 184,276)"
], RGBColor(0x00, 0x7A, 0x33))

# 10. Recommendations
add_content_slide("제작 권장사항 (Action Items)", [
    "## style_profile.json 업데이트 필요 항목",
    "- 리스트형 구조 추가 (30~40초, 번호리스트, 이미지 비율 높음)",
    "- 3자 시점 구조 추가 (친구/주변인 이야기 포맷)",
    "- 가짜후기 반박 장면 추가 (선제적 의심 해소)",
    "- 긴급감 CTA 마무리 옵션 추가 (40%가 사용)",
    "",
    "## 훅킹 이미지 제작 우선순위",
    "- 1순위: 네이트판/패완표 스타일 커뮤니티 게시글 스크린샷",
    "- 2순위: 민감 주제용 일러스트/그림 대체",
    "- 3순위: 충격 비주얼 (클로즈업, 팩 사진 등)",
    "",
    "## 제품별 최적 포맷",
    "- 착색크림: 스토리형 50~60초 (공감+좌절 강조)",
    "- 석류정: 다양한 포맷 가능 (스토리/리스트/3자시점 모두)",
    "- 다이어트/체형: 리스트형 30~40초 (실용 정보 중심)",
    "",
    "## 다음 단계",
    "- 프롬프트 v3.0에 새 구조 반영 (수현 담당)",
    "- 훅킹 이미지 템플릿 제작 (민지 담당)"
], RGBColor(0x1A, 0x1A, 0x2E))

prs.save(pptx_path)
print(f"PPT saved: {pptx_path}")
print(f"Total slides: {len(prs.slides)}")
