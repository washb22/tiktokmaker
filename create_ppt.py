# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color palette
BG_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)
BG_LIGHT = RGBColor(0xF5, 0xF5, 0xF5)
ACCENT_PINK = RGBColor(0xFF, 0x6B, 0x9D)
ACCENT_BLUE = RGBColor(0x4E, 0xC5, 0xD1)
ACCENT_PURPLE = RGBColor(0x7B, 0x68, 0xEE)
TEXT_DARK = RGBColor(0x2D, 0x2D, 0x2D)
TEXT_GRAY = RGBColor(0x66, 0x66, 0x66)
TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT_NAME = "Malgun Gothic"

def set_font(run, size=14, bold=False, color=TEXT_DARK, name=FONT_NAME):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, size=14, bold=False, color=TEXT_DARK, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    set_font(run, size, bold, color)
    return txBox

def add_bullet_text(slide, left, top, width, height, items, size=13, color=TEXT_DARK, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = spacing
        run = p.add_run()
        run.text = item
        set_font(run, size, False, color)
    return txBox

# ============================================
# SLIDE 1: Title
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(2.2), Inches(5.333), Pt(4))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_PINK
shape.line.fill.background()

add_text_box(slide, Inches(2), Inches(2.5), Inches(9.333), Inches(1.5),
    "\ud2f1\ud1a1 \uc219\ud3fc \uc601\uc0c1 10\uac1c \ubd84\uc11d \ubcf4\uace0\uc11c", 40, True, TEXT_WHITE, PP_ALIGN.CENTER)

add_text_box(slide, Inches(2), Inches(4.2), Inches(9.333), Inches(0.8),
    "10\ub300 \ubdf0\ud2f0/\ub2e4\uc774\uc5b4\ud2b8 \uce74\ud14c\uace0\ub9ac | \ud504\ub808\uc784 \ub2e8\uc704 \uc815\ubc00 \ubd84\uc11d", 18, False, RGBColor(0xAA, 0xAA, 0xCC), PP_ALIGN.CENTER)

add_text_box(slide, Inches(2), Inches(5.2), Inches(9.333), Inches(0.6),
    "\ubd84\uc11d: \uc7ac\uc6d0(\ub370\uc774\ud130 \ubd84\uc11d\uac00) | \ub9c8\ucf00\ud305: \ubbfc\uc9c0 | \ub300\ubcf8: \uc218\ud604", 14, False, RGBColor(0x88, 0x88, 0xAA), PP_ALIGN.CENTER)

add_text_box(slide, Inches(2), Inches(5.8), Inches(9.333), Inches(0.5),
    "2026.04.01", 13, False, RGBColor(0x77, 0x77, 0x99), PP_ALIGN.CENTER)

# ============================================
# SLIDE 2: Analysis Overview
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_WHITE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
    "\ubd84\uc11d \uac1c\uc694", 30, True, BG_DARK)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.1), Inches(2), Pt(3))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_PINK
shape.line.fill.background()

box = add_shape_bg(slide, Inches(0.6), Inches(1.5), Inches(5.8), Inches(5.5), RGBColor(0xF8, 0xF8, 0xFC))
add_text_box(slide, Inches(0.9), Inches(1.6), Inches(5.4), Inches(0.5),
    "\ubd84\uc11d \ub300\uc0c1 \uc601\uc0c1 10\uac1c", 16, True, ACCENT_PURPLE)

videos = [
    "1. \uaca8\ub4dc\ub791\uc774\ub294 \uc65c \uc2e0\uacbd\uc548\uc500 #\uc774\ubfe0\uc9c0\ub294\ubc95 (7\uc7a5\uba74)",
    "2. \ub2e4\ub9ac \uc774\ubfe0\uc9c0\ub294 \ucd08 \ud604\uc2e4\uc801\uc778 \ubc29\ubc95 #\ub2e4\uc774\uc5b4\ud2b8 (4\uc7a5\uba74)",
    "3. \ubc29\ud559\ub3d9\uc548 \ud798\ub0b4\ubcf4\uc5d0! #\ub2e4\uc774\uc5b4\ud2b8\uc790\uadf9 (6\uc7a5\uba74)",
    "4. \ube14\ub799\ud5e4\ub4dc \ube7c\ub3c4\ube7c\ub3c4 \uacc4\uc18d \uc0dd\uae30\ub294\uc774\uc720 (11\uc7a5\uba74)",
    "5. \uc2e0\uacbd\uc4f0\uc774\ub294\uc0ac\ub78c #\uae4c\ub9dd\uc774 #\uad00\ub9ac (11\uc7a5\uba74)",
    "6. \uc774\uac83\ub9cc \ud574\ub3c4 \uc5c4\uccad \uc774\ubfe0\uc838 #\uc774\ubfe0\uc9c0\ub294\ubc95 (4\uc7a5\uba74)",
    "7. \uc774\ub807\uac8c \ud588\ub2e4\ub294\ub370 #\ubbf8\ubc31 #\ud53c\ubd80\uad00\ub9ac (6\uc7a5\uba74)",
    "8. \uc9c4\uc9dc \uc774\ubfe0\uc838\uc11c\uc634 #\uc774\ubfe0\uc9c0\uc790\uace0 (7\uc7a5\uba74)",
    "9. \ud53c\ubd80 \ud558\uc598\uc9c0\uace0 \uc2f6\uc740 \uc0ac\ub78c\uc740 \uaf2d \ubd10 (4\uc7a5\uba74)",
    "10. \ud53c\uc9c0 \uc5e5\uc5b4\uc9dc\uc9c0\ub9c8 #\ud53c\uc9c0 (6\uc7a5\uba74)",
]
add_bullet_text(slide, Inches(0.9), Inches(2.2), Inches(5.4), Inches(4.5), videos, 12, TEXT_DARK)

box2 = add_shape_bg(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.5), RGBColor(0xF8, 0xF8, 0xFC))
add_text_box(slide, Inches(7.1), Inches(1.6), Inches(5.4), Inches(0.5),
    "\ubd84\uc11d \ubc29\ubc95\ub860", 16, True, ACCENT_PURPLE)

methods = [
    ">> \ud504\ub808\uc784 \ucd94\ucd9c \ubd84\uc11d: \uc601\uc0c1\ubcc4 \uc8fc\uc694 \uc7a5\uba74 \uce90\uccd0 \ud6c4 \uc2dc\uac01\uc801 \ubd84\uc11d",
    ">> \ucd1d \ubd84\uc11d \ud504\ub808\uc784: 66\uac1c \uc7a5\uba74 (\ud3c9\uade0 6.6\uc7a5\uba74/\uc601\uc0c1)",
    ">> \ubd84\uc11d \ud56d\ubaa9:",
    "   - \uccab \uc7a5\uba74 \ud6c5\ud0b9 \uc720\ud615 (\ub124\uc774\ud2b8\ud310/\uc9c8\ubb38/\uc774\ubbf8\uc9c0/\ubc08)",
    "   - \ud14d\uc2a4\ud2b8 \ud45c\uc2dc \ubc29\uc2dd (\uc704\uce58, \ud06c\uae30, \uc0c9\uc0c1)",
    "   - \uc774\ubbf8\uc9c0 \uc0ac\uc6a9 \ud328\ud134 (\ud69f\uc218, \uc704\uce58, \uc720\ud615)",
    "   - \uc7a5\uba74 \uc804\ud658 \ud328\ud134 \ubc0f \uad6c\uc870",
    "   - \ub9d0\ud22c/\ubb38\uccb4 \ud2b9\uc9d5",
    "   - \uac10\uc815 \uace1\uc120 \ud750\ub984",
    "",
    ">> \uce74\ud14c\uace0\ub9ac: \ubdf0\ud2f0(\ud53c\ubd80\uad00\ub9ac/\ubbf8\ubc31) 7\uac1c, \ub2e4\uc774\uc5b4\ud2b8 3\uac1c",
    ">> \ud0c0\uac9f: 10\ub300~20\ub300 \ucd08\ubc18 \uc5ec\uc131",
    ">> \uacf5\ud1b5 \ud2b9\uc9d5: \uc804\ubd80 PPL/\uad11\uace0 \ud3ec\ud568 \ucf58\ud150\uce20",
]
add_bullet_text(slide, Inches(7.1), Inches(2.2), Inches(5.4), Inches(4.5), methods, 12, TEXT_DARK)

# ============================================
# SLIDE 3: Common Structure (7-step)
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_WHITE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
    "\uacf5\ud1b5 \uad6c\uc870 \ubd84\uc11d: 7\ub2e8\uacc4 \uc2a4\ud1a0\ub9ac \ud504\ub808\uc784", 30, True, BG_DARK)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.1), Inches(2), Pt(3))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_PINK
shape.line.fill.background()

steps = [
    ("1\ub2e8\uacc4\n\ud6c5\ud0b9", ACCENT_PINK,
     "\uacf5\uac10\ud615 \uc774\ubbf8\uc9c0 + \uc9c8\ubb38/\uacf5\uac10 \ud14d\uc2a4\ud2b8\n\"\ub2e4\ub4e4 \uc774\uc815\ub3c4 \ube14\ub799\ud5e4\ub4dc\ub294 \uac00\uc9c0\uace0\uc788\uc9c0?\"\n\"\uc5ec\ub7ec\ubd84\uc758 \uc18c\uc74c\uc21c\uc740 \uc5b4\ub5a4 \uc0c9\uc0c1\uc785\ub2c8\uae4c?\""),
    ("2\ub2e8\uacc4\n\uacf5\uac10 \ud655\uc7a5", RGBColor(0xFF, 0x8C, 0x69),
     "\uace0\ubbfc \uc2ec\ud654, \uac10\uc815 \uc790\uadf9\n\"\ub098\uc911\uc5d0 \uad00\ub9ac\ud558\uae30 \uc9c4\uc9dc \ube61\uc480...\"\n\"\ubc18\uc5d0\uc11c \ud558\uc580 \uc560\ub4e4\uc774 \ub108\ubb34 \ubd80\ub7ec\uc6e0\uc74c\""),
    ("3\ub2e8\uacc4\n\uc6d0\uc778 \uc124\uba85", RGBColor(0xFF, 0xA5, 0x00),
     "\uc804\ubb38\uc801 \uc9c0\uc2dd \uc804\ub2ec (\ud14d\uc2a4\ud2b8 only)\n\"\ube14\ub799\ud5e4\ub4dc\ub294 \ud53c\uc9c0\uac00 \ubaa8\uacf5\uc18d\uc5d0 \uc313\uc5ec\uc11c\"\n\"\uba5c\ub77c\ub2cc \uc0c9\uc18c\uac00 \ub9ce\uc544\uc11c\""),
    ("4\ub2e8\uacc4\n\uae30\uc874 \ubc29\ubc95\n\ubd80\uc815", RGBColor(0xDD, 0xA0, 0xDD),
     "\uae30\uc874 \ud574\uacb0\ucc45\uc758 \ud55c\uacc4 \uc81c\uc2dc\n\"\ucf54\ud329\uc740 \uc5b5\uc9c0\ub85c \ube14\ub799\ud5e4\ub4dc\ub97c \uc9dc\ub0b4\ub294\uac83\"\n\"\ubcd1\uc6d0\ub3c4 \uc54c\uc544\ubd24\ub294\ub370 \ubcc4 \uc18c\uc6a9 \uc5c6\ub354\ub77c\uad6c\uc694\""),
    ("5\ub2e8\uacc4\n\ud574\uacb0\ucc45\n(\uc81c\ud488\uc18c\uac1c)", ACCENT_BLUE,
     "\uc790\uc5f0\uc2a4\ub7ec\uc6b4 \uc81c\ud488 \ub3c4\uc785\n\"\uadf8\ub798\uc11c \uc81c\uac00 \ucc3e\uc740\uac8c \ucc29\uc0c9\uc804\uc6a9\ud06c\ub9bc\"\n\"\uc6cc\uc2dc\ube44 \ube14\ub799\ud5e4\ub4dc \uc624\uc77c \uc4f0\uad6c\uc788\uc5b4!\""),
    ("6\ub2e8\uacc4\n\uc0ac\uc6a9\ud6c4\uae30\n/\uc99d\uac70", RGBColor(0x66, 0xCD, 0xAA),
     "\ube44\ud3ec/\uc560\ud504\ud130 + \uccb4\ud5d8 \ud6c4\uae30\n\"\ud55c\ub2ec\uc815\ub3c4\uba74 \uafb8\uc900\ud788 \uc368\uc8fc\uba74\"\n\"2\uc8fc\ub9cc\uc5d0 \ub2e4\ub9ac \uc587\uc544\uc9c4\uc36c\u314b\u314b\u314b\""),
    ("7\ub2e8\uacc4\nCTA/\ub9c8\ubb34\ub9ac", RGBColor(0x77, 0x88, 0xEE),
     "\ud589\ub3d9 \ucd09\uad6c + \ubd80\uac00 \ud301\n\"\uc9c0\uae08 \ub2f9\uc7a5 \ubdf0\ud2f0\ub77c\ub3c4 \uad00\ub9ac\ud574!\"\n\"\uc544\ubb34\uac70\ub098 \uc0ac\uc9c0\ub9d0\uace0 \ub3c8\ubc84\ub9ac\uc9c0\ub9c8 \u3160\u3160\u3160\""),
]

for i, (label, color, desc) in enumerate(steps):
    x = Inches(0.5) + Inches(i * 1.8)
    s = add_shape_bg(slide, x, Inches(1.6), Inches(1.6), Inches(1.2), color)
    tf = s.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    set_font(run, 12, True, TEXT_WHITE)
    tf.paragraphs[0].space_before = Pt(8)

    if i < 6:
        arrow_x = x + Inches(1.6)
        add_text_box(slide, arrow_x, Inches(1.85), Inches(0.2), Inches(0.5), ">", 16, True, TEXT_GRAY, PP_ALIGN.CENTER)

    add_text_box(slide, x, Inches(3.0), Inches(1.65), Inches(3.5), desc, 9, False, TEXT_DARK)

box = add_shape_bg(slide, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.9), RGBColor(0xFD, 0xF0, 0xF5))
add_text_box(slide, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.7),
    "\ud575\uc2ec \ubc1c\uacac: 10\uac1c \uc601\uc0c1 \ubaa8\ub450 \"\ud6c5\ud0b9 -> \uacf5\uac10 -> \uc6d0\uc778 -> \ubd80\uc815 -> \uc81c\ud488 -> \ud6c4\uae30 -> CTA\" 7\ub2e8\uacc4 \uad6c\uc870\ub97c \ub530\ub984. \uad11\uace0\uc784\uc5d0\ub3c4 \ubd88\uad6c\ud558\uace0 \uc790\uc5f0\uc2a4\ub7ec\uc6b4 \uc2a4\ud1a0\ub9ac\ud154\ub9c1\uc73c\ub85c \uc81c\ud488\uc744 \ub179\uc5ec\ub0b4\ub294 \ud328\ud134.",
    12, False, ACCENT_PINK)

# ============================================
# SLIDE 4: Text/Visual Style
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_WHITE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
    "\ud14d\uc2a4\ud2b8 / \ube44\uc8fc\uc5bc \uc2a4\ud0c0\uc77c \ubd84\uc11d", 30, True, BG_DARK)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.1), Inches(2), Pt(3))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_BLUE
shape.line.fill.background()

box1 = add_shape_bg(slide, Inches(0.5), Inches(1.5), Inches(3.8), Inches(5.5), RGBColor(0xF0, 0xF8, 0xFF))
add_text_box(slide, Inches(0.7), Inches(1.6), Inches(3.4), Inches(0.4), "\ud14d\uc2a4\ud2b8 \uc2a4\ud0c0\uc77c", 16, True, ACCENT_BLUE)
items1 = [
    ">> \ud3f0\ud2b8: \uace0\ub515 \uacc4\uc5f4, \ub465\uadfc \uc0b0\uc138\ub9ac\ud504 (\ub9d1\uc740\uace0\ub515 \uc720\uc0ac)",
    ">> \ud06c\uae30: \ubcf8\ubb38 16~20pt \ucd94\uc815, \uac15\uc870\ubd80 22~28pt",
    ">> \uc0c9\uc0c1: \uac80\uc815(#2D2D2D) \uae30\ubcf8, \uac15\uc870\uc2dc \ubcfc\ub4dc\uccb4 \ud65c\uc6a9",
    ">> \ud2b9\uc218 \uac15\uc870: \ubc11\uc904/\ubcfc\ub4dc \ud63c\uc6a9, \uc0c9 \ubcc0\uacbd\uc740 \ub4dc\ubb3c",
    ">> \uc904\uac04\uaca9: \ub113\uc74c (1.8~2.0\ubc30), \uac00\ub3c5\uc131 \uadf9\ub300\ud654",
    ">> \uc815\ub82c: \uac00\uc6b4\ub370 \uc815\ub82c \uae30\ubcf8, \uae34 \ud14d\uc2a4\ud2b8\ub294 \uc88c\uce21",
    ">> \uc704\uce58: \uc774\ubbf8\uc9c0 \uc544\ub798\uc5d0 \ud14d\uc2a4\ud2b8 \ubc30\uce58\uac00 \uc8fc\ub958",
    ">> \ud55c \ud654\uba74 \ud14d\uc2a4\ud2b8\uc591: 3~6\uc904\uc774 \ucd5c\uc801",
]
add_bullet_text(slide, Inches(0.7), Inches(2.1), Inches(3.4), Inches(4.5), items1, 11, TEXT_DARK)

box2 = add_shape_bg(slide, Inches(4.6), Inches(1.5), Inches(3.8), Inches(5.5), RGBColor(0xFD, 0xF5, 0xF8))
add_text_box(slide, Inches(4.8), Inches(1.6), Inches(3.4), Inches(0.4), "\ubc30\uacbd \uc2a4\ud0c0\uc77c", 16, True, ACCENT_PINK)
items2 = [
    ">> \uae30\ubcf8 \ubc30\uacbd: \uc21c\ubc31(#FFFFFF) \ub610\ub294 \uc5f0\ud68c\uc0c9(#F5F5F5)",
    ">> \uc774\ubbf8\uc9c0 \uc7a5\uba74: \uc0c1\ub2e8 \uc774\ubbf8\uc9c0 + \ud558\ub2e8 \ud770\ubc30\uacbd \ud14d\uc2a4\ud2b8",
    ">> \ud14d\uc2a4\ud2b8 \uc804\uc6a9: \uc644\uc804 \ud770\ubc30\uacbd\uc5d0 \ud14d\uc2a4\ud2b8\ub9cc",
    ">> \ubc08/\uce90\uccd0: \uc720\ud29c\ube0c/\uc608\ub2a5 \uce90\uccd0 + \uc790\ub9c9 \uc624\ubc84\ub808\uc774",
    ">> \uc81c\ud488 \uc7a5\uba74: \uc2e4\uc81c \ucd2c\uc601 \uc0ac\uc9c4 (\uc0dd\ud65c\uac10 \uc788\ub294 \ubc30\uacbd)",
    ">> \uc0c9\uc0c1 \ubc30\uacbd: \uc81c\ud488 \uc124\uba85\uc2dc \ud30c\uc2a4\ud154\ud1a4 \uac00\ub054 \uc0ac\uc6a9",
    "   (\uc608: \uc5f0\ud551\ud06c, \uc5f0\ud30c\ub791 - \ud654\uc7a5\ud488 \uc131\ubd84 \uc124\uba85)",
    ">> \uc804\uccb4\uc801\uc73c\ub85c \uae54\ub054\ud558\uace0 \ubbf8\ub2c8\uba40\ud55c \ub290\ub08c \uc720\uc9c0",
]
add_bullet_text(slide, Inches(4.8), Inches(2.1), Inches(3.4), Inches(4.5), items2, 11, TEXT_DARK)

box3 = add_shape_bg(slide, Inches(8.7), Inches(1.5), Inches(4.1), Inches(5.5), RGBColor(0xF0, 0xFF, 0xF0))
add_text_box(slide, Inches(8.9), Inches(1.6), Inches(3.7), Inches(0.4), "\ub808\uc774\uc544\uc6c3 \ud328\ud134", 16, True, RGBColor(0x2E, 0x8B, 0x57))
items3 = [
    ">> \ud328\ud134 A (60%): \uc0c1\ub2e8 \uc774\ubbf8\uc9c0 + \ud558\ub2e8 \ud14d\uc2a4\ud2b8",
    "   - \uc774\ubbf8\uc9c0\uac00 \ud654\uba74 \uc0c1\ub2e8 40~50% \ucc28\uc9c0",
    "   - \ud14d\uc2a4\ud2b8\uac00 \ud558\ub2e8\uc5d0 \uac00\uc6b4\ub370 \uc815\ub82c",
    "",
    ">> \ud328\ud134 B (25%): \ud14d\uc2a4\ud2b8 \uc804\uc6a9 \ud654\uba74",
    "   - \ud770\ubc30\uacbd\uc5d0 \ud070 \ud14d\uc2a4\ud2b8",
    "   - \uc124\uba85/\uc6d0\uc778 \ub2e8\uacc4\uc5d0\uc11c \uc8fc\ub85c \uc0ac\uc6a9",
    "",
    ">> \ud328\ud134 C (15%): \ubc08/\uce90\uccd0 + \uc790\ub9c9",
    "   - \uc608\ub2a5/\uc720\ud29c\ube0c \uce90\uccd0 \uc704\uc5d0 \uc790\ub9c9 \uc624\ubc84\ub808\uc774",
    '   - "\uc5b4\uc6b0 \ub118\uc5b4\uac08\uac8c\uc694", "\uc644\ubcbd\ud558\ub2e4!" \ub4f1 \ub9ac\uc561\uc158',
    "",
    ">> \uc138\ub85c\ud615(9:16) \ucd5c\uc801\ud654, \uc5ec\ubc31 \ucda9\ubd84\ud788 \ud655\ubcf4",
]
add_bullet_text(slide, Inches(8.9), Inches(2.1), Inches(3.7), Inches(4.5), items3, 11, TEXT_DARK)

# ============================================
# SLIDE 5: Image Usage Pattern
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_WHITE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
    "\uc774\ubbf8\uc9c0 \uc0ac\uc6a9 \ud328\ud134 \ubd84\uc11d", 30, True, BG_DARK)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.1), Inches(2), Pt(3))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_PURPLE
shape.line.fill.background()

add_text_box(slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.5),
    "\uc601\uc0c1\ubcc4 \uc774\ubbf8\uc9c0 \uc0ac\uc6a9 \ud604\ud669", 16, True, ACCENT_PURPLE)

data = [
    ["\uc601\uc0c1", "\ucd1d\uc7a5\uba74", "\uc774\ubbf8\uc9c0\uc7a5\uba74", "\ube44\uc728", "\uc774\ubbf8\uc9c0\uc720\ud615", "\uccab\uc7a5\uba74"],
    ["\uaca8\ub4dc\ub791\uc774\ub294 \uc65c \uc2e0\uacbd\uc548\uc500", "7", "4", "57%", "\uc778\ubb3c\uc0ac\uc9c4, \ud53c\ubd80\ud074\ub85c\uc988\uc5c5, \uc81c\ud488\uc0ac\uc9c4", "\uc778\ubb3c \uc774\ubbf8\uc9c0+\uacf5\uac10\ud14d\uc2a4\ud2b8"],
    ["\ub2e4\ub9ac \uc774\ubfe0\uc9c0\ub294 \ubc29\ubc95", "4", "4", "100%", "\ube44\ud3ec\uc560\ud504\ud130 \uc0ac\uc9c4, \uc81c\ud488\uc0ac\uc9c4, \ucc29\uc6a9\uc0f7", "\ube44\ud3ec/\uc560\ud504\ud130 \uc774\ubbf8\uc9c0"],
    ["\ubc29\ud559\ub3d9\uc548 \ud798\ub0b4\ubcf4\uc5d0", "6", "4", "67%", "\uc608\ub2a5\uce90\uccd0, \uc77c\ub7ec\uc2a4\ud2b8, \uad11\uace0\uce90\uccd0, \ub3d9\ubb3c\ubc08", "\uc608\ub2a5 \ubc08(\ub0a8\uc790 \uc720\ud29c\ubc84)"],
    ["\ube14\ub799\ud5e4\ub4dc \uacc4\uc18d \uc0dd\uae30\ub294\uc774\uc720", "11", "4", "36%", "\ucf54\ud074\ub85c\uc988\uc5c5, \uc131\ubd84\ub3c4\ud45c", "\ucf54 \ud074\ub85c\uc988\uc5c5 \uc0ac\uc9c4"],
    ["\uc2e0\uacbd\uc4f0\uc774\ub294\uc0ac\ub78c #\uae4c\ub9dd\uc774", "11", "3", "27%", "\uc0c9\uc0c1\ucc28\ud2b8, \uc758\uc0ac\uc0ac\uc9c4, \ud53c\ubd80\uc804\ud6c4", "\ud034\uc988\ud615 \uc0c9\uc0c1\ucc28\ud2b8"],
    ["\uc774\uac83\ub9cc \ud574\ub3c4 \uc774\ubfe0\uc838", "4", "4", "100%", "\ub2e4\ub9ac\ube44\ud3ec\uc560\ud504\ud130, \uc81c\ud488, \uc57c\uc678\uc0ac\uc9c4", "\ub2e4\ub9ac \ube44\ud3ec \uc0ac\uc9c4+\uacb0\uacfc\ud14d\uc2a4\ud2b8"],
    ["\uc774\ub807\uac8c \ud588\ub2e4\ub294\ub370 #\ubbf8\ubc31", "6", "4", "67%", "\uc778\ubb3c\uc0ac\uc9c4, \uc81c\ud488\uc0ac\uc9c4, \uc190\ud074\ub85c\uc988\uc5c5", "\uc778\ubb3c(\uce5c\uad6c) \uc0ac\uc9c4+\uc2a4\ud1a0\ub9ac"],
    ["\uc9c4\uc9dc \uc774\ubfe0\uc838\uc11c\uc634", "7", "5", "71%", "\uc608\ub2a5\uce90\uccd0, \uc74c\uc2dd\uc0ac\uc9c4, \uc81c\ud488\uc0ac\uc9c4", "\uc608\ub2a5\uce90\uccd0(\uc2dc\ud5d8\uc885\ub8cc \ubc08)"],
    ["\ud53c\ubd80 \ud558\uc598\uc9c0\uace0 \uc2f6\uc740 \uc0ac\ub78c", "4", "4", "100%", "\ud53c\ubd80\ud074\ub85c\uc988\uc5c5, \uc608\ub2a5\uce90\uccd0, \ud53c\ubd80\uc0ac\uc9c4", "\uc5b4\ub450\uc6b4\ud53c\ubd80 \ud074\ub85c\uc988\uc5c5+\uacf5\uac10"],
    ["\ud53c\uc9c0 \uc5e5\uc5b4\uc9dc\uc9c0\ub9c8", "6", "5", "83%", "\ud53c\uc9c0\ud074\ub85c\uc988\uc5c5, \uc5bc\uad74\uc0ac\uc9c4, \uc81c\ud488\uc0ac\uc9c4, X\ud45c\uc2dc", "\ud53c\uc9c0 \ud074\ub85c\uc988\uc5c5+\uacbd\uace0"],
]

y_start = Inches(2.1)
row_h = Inches(0.45)

for row_idx, row in enumerate(data):
    y = y_start + row_h * row_idx
    col_widths = [Inches(2.5), Inches(0.8), Inches(1.0), Inches(0.7), Inches(4.3), Inches(3.0)]
    x = Inches(0.3)
    for col_idx, cell in enumerate(row):
        bg_color = RGBColor(0xE8, 0xE8, 0xF0) if row_idx == 0 else (RGBColor(0xFA, 0xFA, 0xFF) if row_idx % 2 == 0 else BG_WHITE)
        s = add_shape_bg(slide, x, y, col_widths[col_idx], row_h, bg_color)
        s.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
        s.line.width = Pt(0.5)

        tf = s.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER if col_idx in [1,2,3] else PP_ALIGN.LEFT
        run = tf.paragraphs[0].add_run()
        run.text = cell
        font_size = 9 if row_idx > 0 else 10
        set_font(run, font_size, row_idx == 0, TEXT_DARK if row_idx > 0 else RGBColor(0x3A, 0x3A, 0x5C))
        x += col_widths[col_idx]

box = add_shape_bg(slide, Inches(0.3), Inches(6.3), Inches(12.6), Inches(0.8), RGBColor(0xF5, 0xF0, 0xFF))
items_sum = [
    "\ud3c9\uade0 \uc774\ubbf8\uc9c0 \ube44\uc728: 70.8% | \uc774\ubbf8\uc9c0 \uc720\ud615 TOP 3: 1) \uc81c\ud488 \uc2e4\uc0ac\uc9c4 2) \ud53c\ubd80/\uc2e0\uccb4 \ud074\ub85c\uc988\uc5c5 3) \uc608\ub2a5/\ubc08 \uce90\uccd0 | \uccab\uc7a5\uba74 \uc774\ubbf8\uc9c0 \uc0ac\uc6a9\ub960: 100% (10/10\uac1c \uc601\uc0c1)",
]
add_bullet_text(slide, Inches(0.6), Inches(6.35), Inches(12.0), Inches(0.7), items_sum, 12, ACCENT_PURPLE)

# ============================================
# SLIDE 6: Tone/Style Analysis
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_WHITE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
    "\ub9d0\ud22c / \ubb38\uccb4 \ubd84\uc11d", 30, True, BG_DARK)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.1), Inches(2), Pt(3))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_PINK
shape.line.fill.background()

box1 = add_shape_bg(slide, Inches(0.5), Inches(1.5), Inches(6), Inches(2.5), RGBColor(0xFD, 0xF5, 0xF8))
add_text_box(slide, Inches(0.7), Inches(1.6), Inches(5.6), Inches(0.4), "\ud575\uc2ec \ub9d0\ud22c \ud2b9\uc9d5", 16, True, ACCENT_PINK)
tone_items = [
    '>> \ubc18\ub9d0 \uae30\ubcf8: "~\uac70\ub4e0?", "~\uc796\uc544", "~\ub354\ub77c\uad6c", "~\ud560\uac78?"',
    '>> 10\ub300 \uc778\ud130\ub137 \uc6a9\uc5b4: "\u314b\u314b\u314b", "\u3160\u3160", "\u3160\u3160", "\u3147\u3147", "\ub118", "\uc9c4\uc9dc"',
    '>> \uce5c\ud55c \uce5c\uad6c\uc5d0\uac8c \ub9d0\ud558\ub294 \ub4ef\ud55c \ud1a4: "\ub0b4\uac00 \uc54c\ub824\uc904\uac8c", "\ub098\ub3c4 \uadf8\ub7ac\ub294\ub370"',
    '>> \ubb38\uc7a5 \ub05d \uc774\ubaa8\ud2f0\ucf58/\ud2b9\uc218\uae30\ud638 \ud65c\uc6a9: "~\uc36c\u314b\u314b", "~\ube61\uc480...", "~\u3160\u3160\u3160"',
    '>> \uacfc\uc7a5\ud615 \ud45c\ud604 \ube48\ubc88: "\uc9c4\uc9dc \ubbf8\ucce4\uc5b4", "\uc5c4\uccad \uc774\ubfe0\uc838", "\uc644\uc804 \ub2ec\ub77c\uc9d0"',
]
add_bullet_text(slide, Inches(0.7), Inches(2.1), Inches(5.6), Inches(1.8), tone_items, 11, TEXT_DARK)

box2 = add_shape_bg(slide, Inches(6.8), Inches(1.5), Inches(6), Inches(2.5), RGBColor(0xF0, 0xF8, 0xFF))
add_text_box(slide, Inches(7.0), Inches(1.6), Inches(5.6), Inches(0.4), "\uc124\ub4dd \ud654\ubc95 \ud328\ud134", 16, True, ACCENT_BLUE)
persuade = [
    '>> \uacbd\ud5d8\ub2f4 \ud654\ubc95: "\ub098\ub3c4 \uc6d0\ub798~", "\ub0b4 \uce5c\uad6c\uac00~", "\uc6b0\ub9ac\ubc18 \uc560\uac00~"',
    '>> \ubd80\uc815 \uac15\uc870: "\uc808\ub300 \uc548\ub3fc", "\uc5e5\uc5b4\uc9dc\uc9c0\ub9c8", "\ub3c8\ubc84\ub9ac\uc9c0\ub9c8"',
    '>> \uc2dc\uac04 \ud55c\uc815: "2\uc8fc\ub9cc\uc5d0", "\ud55c\ub2ec\uc774\uba74", "\ubc29\ud559\ub3d9\uc548\ub9cc"',
    '>> \uc0ac\ud68c\uc801 \uc99d\uac70: "\ub2e4\ub4e4 \ud558\ub098\uc529 \uc788\uc796\uc544", "\uce5c\uad6c\uac00 \ucd94\ucc9c\ud574\uc918\uc11c"',
    '>> \uacf5\ud3ec/\ubd88\uc548 \uc790\uadf9: "\ub098\uc911\uc5d0 \ub0a8\uc790\uce5c\uad6c \uc0dd\uae30\uae30\ub3c4 \ud560\uac74\ub370"',
]
add_bullet_text(slide, Inches(7.0), Inches(2.1), Inches(5.6), Inches(1.8), persuade, 11, TEXT_DARK)

box3 = add_shape_bg(slide, Inches(0.5), Inches(4.3), Inches(12.3), Inches(2.8), RGBColor(0xF5, 0xF5, 0xFA))
add_text_box(slide, Inches(0.7), Inches(4.4), Inches(11.9), Inches(0.4), "\uc601\uc0c1\ubcc4 \ub300\ud45c \ubb38\uccb4 \uc608\uc2dc", 16, True, ACCENT_PURPLE)

examples = [
    '"\uc694\uc998 \ud559\uc0dd\uc560\ub4e4 \ubcf4\uba74 \uc804\ubd80 \uc678\ubaa8\ub9cc \uc2e0\uacbd\uc4f0\ub354\ub77c" - \uacf5\uac10\ud615 \uc9c4\uc785, \ub610\ub798 \uad00\ucc30\uc790 \uc2dc\uc810',
    '"\ub2e4\uc774\uc5b4\ud2b8 \uc5c6\uc774 \uc6b4\ub3d9 \uc5c6\uc774 \ud5c8\ubc85\uc9c0 \uc587\uc544\uc9c0\ub294 \ubc95 \uc54c\ub824\uc904\uac8c" - \uc194\ub8e8\uc158 \uc57d\uc18d\ud615, \uc870\uac74 \uc5c6\uc74c \uac15\uc870',
    '"\uc5b4\uc6b0 \ub118\uc5b4\uac08\uac8c\uc694... \ubb34\uc790\uc131\uc73c\ub85c \ucc0c\uc9c8\uac70\ub9b4 \uc0ac\ub78c\uc740 \uadf8\ub0e5 \ub118\uae30\uace0" - \ud544\ud130\ub9c1\ud615, \uc9c4\uc131 \ud0c0\uac9f\ub9cc \ub0a8\uae30\uae30',
    '"\ub2e4\ub4e4 \uc774\uc815\ub3c4 \ube14\ub799\ud5e4\ub4dc\ub294 \uac00\uc9c0\uace0\uc788\uc9c0?" - \uc9c8\ubb38\ud615, \uacf5\uac10 \uc720\ub3c4 + \uad00\uc2ec \ub04c\uae30',
    '"\ud639\uc2dc \uc18c\uc911\uc774\uac00 \uae4c\ub9e4\uc11c \uac71\uc815\uc774\uc2e0\ubd84 \uc788\ub098\uc694?" - \ubbfc\uac10\ud55c \uc8fc\uc81c \ubd80\ub4dc\ub7fd\uac8c \uc811\uadfc',
    '"\ucf54\uc5d0 \uc624\ub3cc\ud1a0\ub3cc\ud55c \ud53c\uc9c0\ub098 \ube14\ub799\ud5e4\ub4dc \uac19\uc740\uac70 \uc190\uc73c\ub85c \uc5e5\uc5b4 \ube7c\uc9c0\ub294 \uc9d3 \uc81c\ubc1c \ud558\uc9c0\ub9c8" - \uacbd\uace0\ud615, \uac15\ud55c \uac10\uc815',
    '"\uc9c0\ub09c\uc8fc\uc5d0 \uc2dc\ud5d8 \ub05d\ub0ac\ub294\ub370... \ud55c \uc560\uac00 \uc0b4\uc774 \ud655 \ube60\uc9c4\uac70\uc57c" - \uc2a4\ud1a0\ub9ac\ud154\ub9c1\ud615, \uc8fc\ubcc0\uc778\ubb3c \ub3c4\uc785',
]
add_bullet_text(slide, Inches(0.7), Inches(4.9), Inches(11.9), Inches(2.0), examples, 11, TEXT_DARK, Pt(4))

# ============================================
# SLIDE 7: Emotion Curve
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_WHITE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
    "\uac10\uc815 \uace1\uc120 \ubd84\uc11d", 30, True, BG_DARK)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.1), Inches(2), Pt(3))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xFF, 0xA5, 0x00)
shape.line.fill.background()

stages = [
    ("\ud6c5\ud0b9\n(\ud638\uae30\uc2ec)", 6, ACCENT_PINK, "\uad00\uc2ec \uc720\ubc1c\n\"\ub098\ub3c4 \ud574\ub2f9\ub418\ub098?\""),
    ("\uacf5\uac10\n(\ubd88\uc548)", 8, RGBColor(0xFF, 0x45, 0x00), "\uac10\uc815 \uace0\uc870\n\"\ub9de\uc544 \ub098\ub3c4 \uadf8\ub798\""),
    ("\uc6d0\uc778\n(\uc774\ud574)", 5, RGBColor(0xFF, 0xA5, 0x00), "\uc9c0\uc2dd \uc2b5\ub4dd\n\"\uc544 \uadf8\ub798\uc11c\uad6c\ub098\""),
    ("\ubd80\uc815\n(\uc88c\uc808)", 9, RGBColor(0xDC, 0x14, 0x3C), "\uac10\uc815 \ucd5c\uace0\uc810\n\"\uadf8\ub7fc \uc5b4\uca4c\uc9c0...\""),
    ("\ud574\uacb0\ucc45\n(\uae30\ub300)", 7, ACCENT_BLUE, "\ud76c\ub9dd \uc804\ud658\n\"\uc774\uac8c \uc788\uc5c8\ub124!\""),
    ("\ud6c4\uae30\n(\uc548\uc2ec)", 6, RGBColor(0x32, 0xCD, 0x32), "\uc2e0\ub8b0 \ud655\ubcf4\n\"\ud6a8\uacfc \uc788\uad6c\ub098\""),
    ("CTA\n(\ud589\ub3d9)", 7, ACCENT_PURPLE, "\ud589\ub3d9 \uc720\ub3c4\n\"\ub098\ub3c4 \ud574\ubd10\uc57c\uc9c0\""),
]

for i, (label, intensity, color, desc) in enumerate(stages):
    x = Inches(0.5) + Inches(i * 1.8)
    bar_height = Inches(intensity * 0.35)
    bar_y = Inches(5.5) - bar_height

    s = add_shape_bg(slide, x + Inches(0.2), bar_y, Inches(1.2), bar_height, color)
    tf = s.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    set_font(run, 10, True, TEXT_WHITE)

    add_text_box(slide, x, Inches(5.6), Inches(1.7), Inches(1.2), desc, 9, False, TEXT_DARK, PP_ALIGN.CENTER)

box = add_shape_bg(slide, Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.5), RGBColor(0xFF, 0xF8, 0xE1))
add_text_box(slide, Inches(0.8), Inches(6.85), Inches(11.7), Inches(0.4),
    "\ud575\uc2ec: \"\ubd80\uc815\" \ub2e8\uacc4\uc5d0\uc11c \uac10\uc815\uc744 \ucd5c\uace0\ub85c \ub04c\uc5b4\uc62c\ub9b0 \ub4a4, \"\ud574\uacb0\ucc45\"\uc73c\ub85c \uae09\uc804\ud658\ud558\ub294 \uac83\uc774 \uad6c\ub9e4 \uc804\ud658\uc758 \ud575\uc2ec \uba54\ucee4\ub2c8\uc998. 4\ub2e8\uacc4->5\ub2e8\uacc4 \uc804\ud658\uc774 \uac00\uc7a5 \uc911\uc694.",
    12, False, RGBColor(0xCC, 0x7A, 0x00))

# ============================================
# SLIDE 8: 10-teen Target Strategy (Minji)
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_WHITE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
    "10\ub300 \ud0c0\uac9f \uc804\ub7b5 \ubd84\uc11d (\ubbfc\uc9c0 \ub9c8\ucf00\ud130 \uad00\uc810)", 30, True, BG_DARK)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.1), Inches(2), Pt(3))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_PINK
shape.line.fill.background()

box1 = add_shape_bg(slide, Inches(0.5), Inches(1.5), Inches(4), Inches(2.5), RGBColor(0xFD, 0xF5, 0xF8))
add_text_box(slide, Inches(0.7), Inches(1.6), Inches(3.6), Inches(0.4), "10\ub300 \uc2ec\ub9ac \uacf5\ub7b5 \ud3ec\uc778\ud2b8", 15, True, ACCENT_PINK)
psych = [
    ">> \uc678\ubaa8 \ucf64\ud50c\ub809\uc2a4 \uadf9\ub300\ud654 \uc2dc\uae30",
    '   - "\ubc18 \uc560\ub4e4 \ub2e4 \uc774\ubfe0", "\ub0a8\uc790\uce5c\uad6c \uc0dd\uae38\ub54c"',
    ">> \ub610\ub798 \ube44\uad50 \uc2ec\ub9ac \ud65c\uc6a9",
    '   - "\uc6b0\ub9ac\ubc18 \ud53c\ubd80\ud1a4\uc774", "\uce5c\uad6c\uac00 \uc54c\ub824\uc918\uc11c"',
    ">> \ubc29\ud559/\uc2dc\ud5d8 \ub4f1 \ud559\uad50 \uc0ac\uc774\ud074 \uc5f0\uacc4",
    '   - "\ubc29\ud559\ub3d9\uc548 \ud798\ub0b4\ubcf4\uc5d0", "\uc2dc\ud5d8 \ub05d\ub098\uace0"',
    ">> \uc6a9\ub3c8 \ubc94\uc704 \ub0b4 \uc81c\ud488 \ucd94\ucc9c",
    '   - \uac00\uaca9 \uc9c1\uc811 \uc5b8\uae09 \ud68c\ud53c, "\ubd80\ub2f4\uc5c6\uc774" \uac15\uc870',
]
add_bullet_text(slide, Inches(0.7), Inches(2.1), Inches(3.6), Inches(1.8), psych, 11, TEXT_DARK)

box2 = add_shape_bg(slide, Inches(4.8), Inches(1.5), Inches(4), Inches(2.5), RGBColor(0xF0, 0xF8, 0xFF))
add_text_box(slide, Inches(5.0), Inches(1.6), Inches(3.6), Inches(0.4), "\ud50c\ub7ab\ud3fc \ucd5c\uc801\ud654 \uc804\ub7b5", 15, True, ACCENT_BLUE)
platform = [
    ">> \uc378\ub124\uc77c = \uccab\ud504\ub808\uc784 = \uac00\uc7a5 \uc911\uc694",
    "   - 10\uac1c \uc911 10\uac1c \ubaa8\ub450 \uc774\ubbf8\uc9c0+\ud14d\uc2a4\ud2b8 \uc870\ud569",
    ">> \ud574\uc2dc\ud0dc\uadf8 \uc804\ub7b5",
    "   - #\uc774\ubfe0\uc9c0\ub294\ubc95 #\uad00\ub9ac #\ud53c\ubd80\uad00\ub9ac \uacc4\uc5f4 \uc8fc\ub825",
    "   - \uac10\uc815\ud615 \ud0dc\uadf8: #\uc774\ubfe0\uc9c0\uc790\uace0 #\ub2e4\uc774\uc5b4\ud2b8\uc790\uadf9",
    ">> \uc601\uc0c1 \uae38\uc774: 4~11\uc7a5\uba74 (30\ucd08~90\ucd08 \ucd94\uc815)",
    ">> \ub313\uae00 \uc720\ub3c4: \uc9c8\ubb38\ud615 \ub9c8\ubb34\ub9ac\ub85c \ucc38\uc5ec \ucd09\uad6c",
    ">> \uc2dc\ub9ac\uc988\ud654: \uc720\uc0ac \uc8fc\uc81c \ubc18\ubcf5\uc73c\ub85c \ud314\ub85c\uc6b0 \uc720\ub3c4",
]
add_bullet_text(slide, Inches(5.0), Inches(2.1), Inches(3.6), Inches(1.8), platform, 11, TEXT_DARK)

box3 = add_shape_bg(slide, Inches(9.1), Inches(1.5), Inches(3.7), Inches(2.5), RGBColor(0xF0, 0xFF, 0xF0))
add_text_box(slide, Inches(9.3), Inches(1.6), Inches(3.3), Inches(0.4), "\ucf58\ud150\uce20 \ubbf9\uc2a4 \ube44\uc728", 15, True, RGBColor(0x2E, 0x8B, 0x57))
mix_items = [
    ">> \ud53c\ubd80\uad00\ub9ac/\ubbf8\ubc31: 50% (5/10)",
    ">> \ub2e4\uc774\uc5b4\ud2b8/\uccb4\ud615: 30% (3/10)",
    ">> \ubc14\ub514\uad00\ub9ac(\uae30\ud0c0): 20% (2/10)",
    "",
    ">> PPL \ud3ec\ud568: 100% (10/10)",
    ">> \ub124\uc774\ud2f0\ube0c \uad11\uace0 \ud615\uc2dd: 100%",
    ">> \uc81c\ud488 \uc9c1\uc811 \ub178\ucd9c: \ud3c9\uade0 2\ud68c/\uc601\uc0c1",
    "",
    ">> \uac00\uc7a5 \ub9ce\uc774 \ub4f1\uc7a5\ud55c \ube0c\ub79c\ub4dc:",
    "   \uc6cc\uc2dc\ube44(WASH.B) - 3\ud68c",
    "   \uc5d0\uc774\uc5b4\ud2b8 - 2\ud68c",
]
add_bullet_text(slide, Inches(9.3), Inches(2.1), Inches(3.3), Inches(1.8), mix_items, 11, TEXT_DARK)

box4 = add_shape_bg(slide, Inches(0.5), Inches(4.3), Inches(12.3), Inches(2.8), RGBColor(0xF8, 0xF0, 0xFF))
add_text_box(slide, Inches(0.7), Inches(4.4), Inches(11.9), Inches(0.4), "\ubbfc\uc9c0\uc758 \ud575\uc2ec \ub9c8\ucf00\ud305 \uc778\uc0ac\uc774\ud2b8", 16, True, ACCENT_PURPLE)
marketing = [
    "1. \"\uad11\uace0\ucc98\ub7fc \ubcf4\uc774\uc9c0 \uc54a\ub294 \uad11\uace0\"\uac00 \ud575\uc2ec - 10\uac1c \uc601\uc0c1 \ubaa8\ub450 \uce5c\uad6c\uac00 \uc54c\ub824\uc8fc\ub294 \ub290\ub08c\uc758 \ub124\uc774\ud2f0\ube0c \uad11\uace0 \ud615\uc2dd",
    "2. \uc81c\ud488 \ub4f1\uc7a5 \ud0c0\uc774\ubc0d\uc774 \uc601\uc0c1 \ud6c4\ubc18 60~70% \uc9c0\uc810 - \ucda9\ubd84\ud55c \uacf5\uac10/\ubb38\uc81c \uc778\uc2dd \ud6c4 \uc790\uc5f0\uc2a4\ub7ec\uc6b4 \ub3c4\uc785",
    "3. \"\ube44\ud3ec/\uc560\ud504\ud130\"\ub294 \ud544\uc218 \uc694\uc18c - 10\uac1c \uc911 8\uac1c\uac00 \uc804\ud6c4 \ube44\uad50 \uc0ac\uc9c4 \ud3ec\ud568, \uc2e0\ub8b0 \uad6c\ucd95\uc758 \ud575\uc2ec \uc7a5\uce58",
    "4. \ubd80\uc815\uc801 \uac10\uc815(\ubd88\uc548/\ucf64\ud50c\ub809\uc2a4) \uc790\uadf9 \ud6c4 \ud574\uacb0\ucc45 \uc81c\uc2dc = \"Pain Point Marketing\"\uc758 \uc804\ud615\uc801 \ud328\ud134",
    "5. \ud0c0\uac9f\uc758 \uc77c\uc0c1 \uc5b8\uc5b4\ub97c \uadf8\ub300\ub85c \uc0ac\uc6a9 - \"\u314b\u314b\", \"\u3160\u3160\", \ubc18\ub9d0\uccb4 \ub4f1\uc73c\ub85c \uad11\uace0 \uac70\ubd80\uac10 \ucd5c\uc18c\ud654",
    "6. 3\uc790 \ucd94\ucc9c \ud654\ubc95(\"\uce5c\uad6c\uac00\", \"\uc6b0\ub9ac\ubc18 \uc560\uac00\") \ud65c\uc6a9\uc73c\ub85c \uac04\uc811 \uacbd\ud5d8 \uc81c\uacf5 - \uc9c1\uc811 \ud310\ub9e4\ubcf4\ub2e4 \ud6a8\uacfc\uc801",
    "7. \ud559\uad50/\uc2dc\ud5d8/\ubc29\ud559 \ub4f1 10\ub300 \ub77c\uc774\ud504\uc0ac\uc774\ud074 \ud0a4\uc6cc\ub4dc \uc801\uadf9 \ud65c\uc6a9 - \uc2dc\uc758\uc131 \ud655\ubcf4",
]
add_bullet_text(slide, Inches(0.7), Inches(4.9), Inches(11.9), Inches(2.0), marketing, 11, TEXT_DARK, Pt(3))

# ============================================
# SLIDE 9: Prompt Design (Suhyun)
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_WHITE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
    "\ud504\ub86c\ud504\ud2b8 \uc124\uacc4 \ud575\uc2ec (\uc218\ud604 \ub300\ubcf8\uc791\uac00 \uad00\uc810)", 30, True, BG_DARK)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.1), Inches(2), Pt(3))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_BLUE
shape.line.fill.background()

box1 = add_shape_bg(slide, Inches(0.5), Inches(1.5), Inches(6), Inches(5.3), RGBColor(0xF0, 0xF8, 0xFF))
add_text_box(slide, Inches(0.7), Inches(1.6), Inches(5.6), Inches(0.4), "\ub300\ubcf8 \ud504\ub86c\ud504\ud2b8 \uad6c\uc870 \ud15c\ud50c\ub9bf", 16, True, ACCENT_BLUE)
template = [
    "[\uc7a5\uba74 1] \ud6c5\ud0b9 (\ud544\uc218: \uc774\ubbf8\uc9c0 + \uc9e7\uc740 \uc9c8\ubb38/\uacf5\uac10 1\uc904)",
    '  - \uc720\ud615 A: \ud074\ub85c\uc988\uc5c5 \uc0ac\uc9c4 + "\ub2e4\ub4e4 ~\ud558\uc796\uc544?"',
    '  - \uc720\ud615 B: \ube44\ud3ec/\uc560\ud504\ud130 + "~\ub9cc\uc5d0 ~\ud55c\uc36c\u314b\u314b"',
    "  - \uc720\ud615 C: \ubc08/\uc608\ub2a5\uce90\uccd0 + \uc0c1\ud669 \uc124\uc815",
    "",
    "[\uc7a5\uba74 2] \uacf5\uac10 \ud655\uc7a5 (\uc774\ubbf8\uc9c0 \uc720\uc9c0 + \ud14d\uc2a4\ud2b8 3~5\uc904 \ucd94\uac00)",
    '  - \ubcf8\uc778 \uacbd\ud5d8 or \uc8fc\ubcc0\uc778 \uacbd\ud5d8 \uc11c\uc220',
    '  - "\ub098\ub3c4 \uc6d0\ub798~", "\uc6b0\ub9ac\ubc18\uc5d0~" \ud328\ud134',
    "",
    "[\uc7a5\uba74 3~4] \uc6d0\uc778/\uc9c0\uc2dd (\ud14d\uc2a4\ud2b8 \uc804\uc6a9 or \ub3c4\ud45c \uc774\ubbf8\uc9c0)",
    "  - \uc804\ubb38 \uc6a9\uc5b4 \uc57d\uac04 \uc11e\ub418 \uc26c\uc6b4 \uc124\uba85",
    '  - "~\uac70\ub4e0?", "~\uc796\uc544" \uc2dd \uc124\uba85\uccb4',
    "",
    "[\uc7a5\uba74 5] \uae30\uc874 \ubc29\ubc95 \ubd80\uc815 (\ud14d\uc2a4\ud2b8 \uc804\uc6a9, \uac10\uc815 \uac15\uc870)",
    '  - "\uadfc\ub370 \uadf8\uac74 ~\ud574\uc11c \uc18c\uc6a9\uc5c6\uc5b4"',
    "  - \ubcfc\ub4dc/\uac15\uc870 \ud14d\uc2a4\ud2b8\ub85c \ud575\uc2ec \ud3ec\uc778\ud2b8",
    "",
    "[\uc7a5\uba74 6] \uc81c\ud488 \ub3c4\uc785 (\uc81c\ud488 \uc2e4\uc0ac\uc9c4 + \uc124\uba85 \ud14d\uc2a4\ud2b8)",
    '  - "\uadf8\ub798\uc11c \ub0b4\uac00 \ucc3e\uc740\uac8c~"',
    "  - \uc790\uc5f0\uc2a4\ub7ec\uc6b4 \uc804\ud658, \uad11\uace0\ud2f0 \ucd5c\uc18c\ud654",
    "",
    "[\uc7a5\uba74 7] \ud6c4\uae30/CTA (\ube44\ud3ec\uc560\ud504\ud130 or \ud14d\uc2a4\ud2b8 \ub9c8\ubb34\ub9ac)",
    "  - \uacb0\uacfc + \ud589\ub3d9 \ucd09\uad6c + \ubd80\uac00 \ud301",
]
add_bullet_text(slide, Inches(0.7), Inches(2.1), Inches(5.6), Inches(4.5), template, 10, TEXT_DARK, Pt(2))

box2 = add_shape_bg(slide, Inches(6.8), Inches(1.5), Inches(6), Inches(2.3), RGBColor(0xFD, 0xF5, 0xF8))
add_text_box(slide, Inches(7.0), Inches(1.6), Inches(5.6), Inches(0.4), "\uc218\ud604\uc758 \ub300\ubcf8 \uc791\uc131 \uaddc\uce59", 16, True, ACCENT_PINK)
rules = [
    "1. \ud55c \uc7a5\uba74\ub2f9 \ud14d\uc2a4\ud2b8: \ucd5c\uc18c 2\uc904, \ucd5c\ub300 7\uc904",
    "2. \ubc18\ub9d0\uccb4 \ud544\uc218, \uc874\ub313\ub9d0\uc740 \uc778\uc6a9\uad6c\uc5d0\uc11c\ub9cc",
    "3. \u314b\u314b, \u3160\u3160 \ub4f1 \uc774\ubaa8\ud2f0\ucf58 \uc7a5\uba74\ub2f9 1~2\uac1c",
    "4. \uc9c8\ubb38\ud615 \ubb38\uc7a5 \uc601\uc0c1\ub2f9 \ucd5c\uc18c 2\uac1c \ud3ec\ud568",
    "5. \uc81c\ud488\uba85 \uc9c1\uc811 \uc5b8\uae09\uc740 \ud6c4\ubc18\ubd80\uc5d0\uc11c\ub9cc",
    '6. "~\ub354\ub77c\uad6c", "~\uac70\ub4e0?" \uc5b4\ubbf8 \uc801\uadf9 \ud65c\uc6a9',
    "7. \uacfc\uc7a5\ud615 \ud45c\ud604\uc740 \ud6c5\ud0b9\uacfc \ud6c4\uae30\uc5d0 \uc9d1\uc911",
]
add_bullet_text(slide, Inches(7.0), Inches(2.1), Inches(5.6), Inches(1.6), rules, 11, TEXT_DARK, Pt(3))

box3 = add_shape_bg(slide, Inches(6.8), Inches(4.0), Inches(6), Inches(2.8), RGBColor(0xF0, 0xFF, 0xF0))
add_text_box(slide, Inches(7.0), Inches(4.1), Inches(5.6), Inches(0.4), "\uc774\ubbf8\uc9c0 \uc9c0\uc2dc\ubb38 \uaddc\uce59", 16, True, RGBColor(0x2E, 0x8B, 0x57))
img_rules = [
    "\uc7a5\uba741: [\uc774\ubbf8\uc9c0: \uc8fc\uc81c \uad00\ub828 \ud074\ub85c\uc988\uc5c5/\ube44\ud3ec\uc0ac\uc9c4]",
    "\uc7a5\uba742: [\uc774\ubbf8\uc9c0: \uc7a5\uba741\uacfc \ub3d9\uc77c or \uad00\ub828 \uc774\ubbf8\uc9c0 \uc720\uc9c0]",
    "\uc7a5\uba743: [\uc774\ubbf8\uc9c0: \uc5c6\uc74c] \ub610\ub294 [\uc774\ubbf8\uc9c0: \ub3c4\ud45c/\uc778\ud3ec\uadf8\ub798\ud53d]",
    "\uc7a5\uba744: [\uc774\ubbf8\uc9c0: \uc5c6\uc74c] \ud14d\uc2a4\ud2b8 \uac15\uc870",
    "\uc7a5\uba745: [\uc774\ubbf8\uc9c0: \ud074\ub85c\uc988\uc5c5/\ud53c\ubd80\uc0ac\uc9c4 + \uac10\uc815\ud14d\uc2a4\ud2b8]",
    "\uc7a5\uba746: [\uc774\ubbf8\uc9c0: \uc81c\ud488 \uc2e4\uc0ac\uc9c4]",
    "\uc7a5\uba747: [\uc774\ubbf8\uc9c0: \ube44\ud3ec/\uc560\ud504\ud130 \ube44\uad50 or \uc5c6\uc74c]",
    "",
    "** \uc774\ubbf8\uc9c0 \uac80\uc0c9 \ud0a4\uc6cc\ub4dc \uc790\ub3d9 \uc0dd\uc131 \ud544\uc218 **",
    "** \uc774\ubbf8\uc9c0 \ube44\uc728: \uc804\uccb4 \uc7a5\uba74\uc758 60~80% **",
]
add_bullet_text(slide, Inches(7.0), Inches(4.6), Inches(5.6), Inches(2.0), img_rules, 11, TEXT_DARK, Pt(3))

# ============================================
# SLIDE 10: Key Insights & Recommendations
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
    "\ud575\uc2ec \uc778\uc0ac\uc774\ud2b8 & \uc81c\uc5b8", 30, True, TEXT_WHITE)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.1), Inches(2), Pt(3))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_PINK
shape.line.fill.background()

insights = [
    ("01", "7\ub2e8\uacc4 \uad6c\uc870\uc758 \ubc95\uce59",
     "10\uac1c \uc601\uc0c1 100%\uac00 \ub3d9\uc77c\ud55c \"\ud6c5\ud0b9-\uacf5\uac10-\uc6d0\uc778-\ubd80\uc815-\ud574\uacb0-\ud6c4\uae30-CTA\" \uad6c\uc870. \uc774 \uad6c\uc870\ub97c \ud504\ub86c\ud504\ud2b8\uc758 \uae30\ubcf8 \uace8\uaca9\uc73c\ub85c \ubc18\ub4dc\uc2dc \uc801\uc6a9\ud574\uc57c \ud568.",
     ACCENT_PINK),
    ("02", "\uc774\ubbf8\uc9c0\uac00 \uc2a4\ud06c\ub864\uc744 \uba48\ucda4\ub2e4",
     "\uccab \uc7a5\uba74 \uc774\ubbf8\uc9c0 \uc0ac\uc6a9\ub960 100%. \ud14d\uc2a4\ud2b8\ub9cc\uc73c\ub85c\ub294 \uc2a4\ud06c\ub864\uc744 \uba48\ucd9c \uc218 \uc5c6\uc74c. \ud504\ub86c\ud504\ud2b8\uc5d0\uc11c \ub9e4 \uc7a5\uba74\uc758 \uc774\ubbf8\uc9c0 \uc9c0\uc2dc\ubb38\uc744 \ud544\uc218 \ucd9c\ub825\ud574\uc57c \ud568.",
     ACCENT_BLUE),
    ("03", "\uad11\uace0\ub97c \uad11\uace0\ucc98\ub7fc \ubcf4\uc774\uc9c0 \uc54a\uac8c",
     "10\uac1c \ubaa8\ub450 PPL\uc774\uc9c0\ub9cc \"\uce5c\uad6c \ucd94\ucc9c\" \ud654\ubc95\uc73c\ub85c \uc704\uc7a5. \ud504\ub86c\ud504\ud2b8\uc5d0\uc11c \"\uc790\uc5f0\uc2a4\ub7ec\uc6b4 \uacbd\ud5d8\ub2f4 \ud1a4\"\uc744 \uba85\uc2dc\uc801\uc73c\ub85c \uc9c0\uc2dc\ud574\uc57c \ud568.",
     ACCENT_PURPLE),
    ("04", "\uac10\uc815 \ub864\ub7ec\ucf54\uc2a4\ud130\uac00 \uc804\ud658\uc728",
     "\"\ubd80\uc815\" \ub2e8\uacc4\uc5d0\uc11c \ubd88\uc548\uc744 \uadf9\ub300\ud654\ud55c \ub4a4 \"\ud574\uacb0\ucc45\"\uc73c\ub85c \uae09\uc804\ud658. \uc774 \uac10\uc815 \ub099\ucc28\uac00 \ud074\uc218\ub85d \uad6c\ub9e4 \uc804\ud658 \ub192\uc74c. \ud504\ub86c\ud504\ud2b8\uc5d0 \uac10\uc815 \uac15\ub3c4 \uc9c0\uc2dc \ud544\uc694.",
     RGBColor(0xFF, 0xA5, 0x00)),
    ("05", "10\ub300 \uc5b8\uc5b4\uc758 \uc815\ud655\ud55c \uc7ac\ud604",
     "\"\u314b\u314b\", \"\u3160\u3160\", \ubc18\ub9d0\uccb4, \"~\uac70\ub4e0?\", \"~\ub354\ub77c\uad6c\" \ub4f1 \ud2b9\uc815 \uc5b4\ubbf8 \ud328\ud134\uc774 \uc2e0\ub8b0\uac10\uc758 \ud575\uc2ec. \ud504\ub86c\ud504\ud2b8\uc5d0 \uc5b4\ubbf8/\ubb38\uccb4 \uc608\uc2dc\ub97c \ubc18\ub4dc\uc2dc \ud3ec\ud568.",
     RGBColor(0x32, 0xCD, 0x32)),
]

for i, (num, title, desc, color) in enumerate(insights):
    y = Inches(1.5) + Inches(i * 1.15)

    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), y, Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = num
    set_font(run, 14, True, TEXT_WHITE)

    add_text_box(slide, Inches(1.6), y, Inches(3), Inches(0.5), title, 16, True, color)
    add_text_box(slide, Inches(1.6), y + Inches(0.4), Inches(10.5), Inches(0.6), desc, 12, False, RGBColor(0xCC, 0xCC, 0xDD))

box = add_shape_bg(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.7), RGBColor(0x2A, 0x2A, 0x4E))
add_text_box(slide, Inches(0.8), Inches(6.55), Inches(11.7), Inches(0.6),
    "\ucd5c\uc885 \uc81c\uc5b8: \ud504\ub86c\ud504\ud2b8\uc5d0 [7\ub2e8\uacc4 \uad6c\uc870 + \uc7a5\uba74\ubcc4 \uc774\ubbf8\uc9c0 \uc9c0\uc2dc + 10\ub300 \ub9d0\ud22c \uc608\uc2dc + \uac10\uc815 \uac15\ub3c4 \ub808\ubca8] 4\uac00\uc9c0\ub97c \ubc18\ub4dc\uc2dc \ud3ec\ud568\uc2dc\ud0ac \uac83. \uc774 4\uc694\uc18c\uac00 \uc219\ud3fc \ub300\ubcf8\uc758 \uc131\ud328\ub97c \uacb0\uc815\ud568.",
    13, True, RGBColor(0xFF, 0xD7, 0x00))

# Save
output_path = r"C:\Users\ghtes\Desktop\tiktok-generator\영상분석_보고서.pptx"
prs.save(output_path)
print(f"PPT saved to: {output_path}")
